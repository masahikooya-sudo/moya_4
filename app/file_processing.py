"""テキストファイル・CSVファイル・Office文書・PDF・JSONのアップロードを受け取り、
マスキング済みの内容を生成するロジック。

表形式データ(CSV/Excel/JSON)は列名(ヘッダー)からPIIの種別を推定して
セル全体を強制マスキングできる(column_matcher.py)。列が明確でない自由形式
データ(txt/docx/pptx/pdf)では、NERによる検出候補を一覧化してユーザーに
確認してもらってからマスキングする2段階フロー(analyze_*/mask_*)を提供する。
"""

import csv
import io
import json
import re

import pymupdf as fitz
import openpyxl
from docx import Document
from presidio_analyzer import RecognizerResult
from pptx import Presentation

from . import column_matcher, config
from .engine import analyze_resolved, mask_full_value, mask_resolved, mask_text

# 日本のビジネス文書で使われやすい文字コードの候補（優先順）。
CANDIDATE_ENCODINGS = ["utf-8-sig", "utf-8", "cp932", "shift_jis", "euc_jp"]

# 分析結果(候補一覧・サンプル値)を画面に表示する際の切り詰め長。
SAMPLE_MAX_LEN = 60

# ラベル直後の値を強制マスキング対象とみなす際の判定パラメータ。
# (PDF/txt等の自由形式文書で「ふりがな」「氏名」等のラベルが単独行として現れ、
#  直後の行に値が続くフォーム的なレイアウトを補完的に検出するために使う)
LABEL_LINE_MAX_LEN = 10
MAX_VALUE_LINES = 3
MAX_VALUE_LINE_LEN = 40


def _open_document(loader, label: str):
    """Office/PDFファイルを開く処理を共通化し、失敗時は ValueError に統一する。

    破損したファイルやパスワード保護されたファイルを開こうとすると、各ライブラリは
    それぞれ異なる例外(RuntimeError, zipfile.BadZipFile, PackageNotFoundError等)を
    送出する。ここで捕捉せずに main.py まで伝播すると、未対応の例外として
    500(素のテキスト応答)になりフロントエンドでJSONパースエラーになってしまうため、
    ValueError に統一し、呼び出し元(main.py)で400として扱えるようにする。
    """
    try:
        return loader()
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(
            f"{label}の読み込みに失敗しました。ファイルが破損しているか、"
            f"対応していない形式・パスワード保護されている可能性があります: {exc}"
        )


def decode_bytes(raw: bytes) -> tuple[str, str]:
    """バイト列を可能な文字コード候補で順にデコードする。

    戻り値: (デコード後文字列, 使用した文字コード)
    """
    last_error = None
    for encoding in CANDIDATE_ENCODINGS:
        try:
            return raw.decode(encoding), encoding
        except (UnicodeDecodeError, LookupError) as exc:
            last_error = exc
    raise ValueError(f"ファイルの文字コードを判定できませんでした: {last_error}")


def _truncate(text: str, limit: int = SAMPLE_MAX_LEN) -> str:
    text = text.replace("\n", " ").replace("\r", " ")
    if len(text) > limit:
        return text[:limit] + "…"
    return text


def _column_entry(key: str, header: str, allowed: set, sample: str) -> dict:
    suggested = column_matcher.match_entity_type(header, allowed)
    return {
        "key": key,
        "header": header,
        "suggested": suggested,
        "suggested_label": config.ENTITY_LABELS_JA.get(suggested) if suggested else None,
        "sample": _truncate(sample or ""),
    }


def _label_forced_results(text: str, allowed_entities: set) -> list:
    """ラベル直後の値を、ラベルに対応するエンティティ種別として強制検出する。

    履歴書等のPDFフォームでは「ふりがな」「氏名」のようなラベルが単独行として
    現れ、直後の行に値が続くレイアウトが多い。こうした短い値(特にふりがなや
    苗字だけの氏名)はNERでは文脈不足のため検出漏れしやすいため、表形式データの
    列名判定(column_matcher.py)と同じキーワード・判定ロジックを流用して補完する。

    ラベル自体が改行で2行に分割されていることがある(例:「ふりが」+「な」)ため、
    1行で一致しなければ次の1行と連結しても試す。値は既定で1行のみを対象とする
    (氏名・ふりがな等は通常1行)が、値の1行目に英数字が含まれる場合はメール
    アドレスや日付のように複数行に分断されている可能性が高いとみなし、
    最大 MAX_VALUE_LINES 行までを1つの値として連結する。
    """
    if not allowed_entities:
        return []

    lines = text.split("\n")
    offsets = []
    pos = 0
    for line in lines:
        offsets.append(pos)
        pos += len(line) + 1

    def label_type_for(line: str):
        stripped = line.strip()
        if not stripped or len(stripped) > LABEL_LINE_MAX_LEN:
            return None
        return column_matcher.match_entity_type(stripped, allowed_entities)

    def exact_label_type_for(line: str):
        stripped = line.strip()
        if not stripped or len(stripped) > LABEL_LINE_MAX_LEN:
            return None
        return column_matcher.match_exact(stripped, allowed_entities)

    results = []
    n = len(lines)
    i = 0
    while i < n:
        # 1行で完全一致すればそれを優先する(例:「氏 名」→氏名)。
        # そうでなければ、次の1行と連結して完全一致するか試す
        # (例:「電話番」+「号」→電話番号)。単独行の部分一致
        # (例:「電話番」が「電話」に部分一致)は、本来2行に分割された
        # ラベルの前半に過ぎない場合があり誤判定しやすいため、
        # 完全一致(単独/連結)がどちらも得られない場合の最終手段とする。
        label_type = exact_label_type_for(lines[i])
        consumed = 1
        if label_type is None and i + 1 < n:
            combined = lines[i].strip() + lines[i + 1].strip()
            # 連結時も完全一致のみを許可する(部分一致だと、無関係な2行を
            # つなげた結果たまたまキーワードを含んでしまうケースを誤判定する)。
            combined_type = exact_label_type_for(combined)
            if combined_type is not None:
                label_type = combined_type
                consumed = 2
        if label_type is None:
            label_type = label_type_for(lines[i])
            consumed = 1

        if label_type is None:
            i += 1
            continue

        j = i + consumed
        value_end = None
        k = j
        taken = 0
        while k < n and taken < MAX_VALUE_LINES:
            candidate = lines[k].strip()
            if not candidate or len(candidate) > MAX_VALUE_LINE_LEN or label_type_for(lines[k]) is not None:
                break
            value_end = k
            taken += 1
            if taken == 1 and not re.search(r"[A-Za-z0-9]", candidate):
                # 氏名・ふりがな等は通常1行の値のため、英数字を含まない場合は
                # ここで打ち切る(次のフィールドのラベルを値として巻き込まないため)。
                break
            k += 1

        if value_end is not None:
            start = offsets[j]
            end = offsets[value_end] + len(lines[value_end])
            span = text[start:end]
            m = re.search(r"\S[\s\S]*\S|\S", span)
            if m:
                results.append(
                    RecognizerResult(
                        entity_type=label_type,
                        start=start + m.start(),
                        end=start + m.end(),
                        score=0.95,
                    )
                )
            i = value_end + 1
        else:
            i = j

    return results


def _analyze_with_label_hints(text: str, entities: list) -> list:
    """NERの検出結果に、ラベル直後の値の強制検出(_label_forced_results)を
    マージして返す。ラベルに基づく検出は確実性が高いため、重複するNER結果より
    優先して採用する。
    """
    allowed = set(entities or config.ALL_ENTITY_CODES)
    ner_results = analyze_resolved(text, entities)
    label_results = _label_forced_results(text, allowed)
    if not label_results:
        return ner_results

    occupied = [(r.start, r.end) for r in label_results]
    merged = list(label_results)
    for result in ner_results:
        if any(result.start < end and result.end > start for start, end in occupied):
            continue
        merged.append(result)
    return sorted(merged, key=lambda r: r.start)


def _collect_candidates(texts, entities) -> list:
    """複数のテキスト断片からNER検出候補を収集し、(種別,一致文字列)で重複排除する。

    一致文字列は前後の空白・改行を取り除いた上で管理する。PDFのページ抽出テキスト
    等では検出スパンの末尾に改行を含むことがあり、マスキング側(mask_text/
    mask_pdf_file)も confirmed の突き合わせ時に同様に strip() するため、
    ここで揃えておかないと確認済み候補が一致せずマスキングされなくなる。
    """
    counts = {}
    order = []
    for text in texts:
        if not text:
            continue
        for result in _analyze_with_label_hints(text, entities):
            value = text[result.start : result.end].strip()
            if not value:
                continue
            key = (result.entity_type, value)
            if key not in counts:
                counts[key] = 0
                order.append(key)
            counts[key] += 1

    return [
        {
            "entity_type": entity_type,
            "entity_label": config.ENTITY_LABELS_JA.get(entity_type, entity_type),
            "text": value,
            "sample": _truncate(value),
            "count": counts[(entity_type, value)],
        }
        for entity_type, value in order
    ]


def _mask_cell(value: str, forced_type, entities: list, style: str):
    """列名からエンティティ種別が確定している場合はセル全体を強制マスキングし、
    そうでなければ通常のNERベースの検出でマスキングする。

    forced_type はユーザーが選択したエンティティ種別(entities)に含まれる場合のみ
    有効とする(画面のチェックボックスを最終的な有効/無効の基準にするため)。
    """
    if forced_type and (not entities or forced_type in entities):
        return mask_full_value(value, forced_type, style)
    return mask_text(value, entities=entities, style=style)


def mask_plain_text_file(raw: bytes, entities: list, style: str, confirmed: set = None):
    """.txt ファイルをマスキングする。"""
    text, _encoding = decode_bytes(raw)
    results = _analyze_with_label_hints(text, entities)
    masked_text, detections = mask_resolved(text, results, style, entities, confirmed)
    return masked_text.encode("utf-8-sig"), detections, text


def analyze_txt_candidates(raw: bytes, entities: list) -> list:
    text, _encoding = decode_bytes(raw)
    return _collect_candidates([text], entities)


def _read_csv_rows(raw: bytes):
    text, _encoding = decode_bytes(raw)
    sample = text[:2048]
    dialect = csv.excel
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;")
    except csv.Error:
        pass
    rows = list(csv.reader(io.StringIO(text), dialect))
    return text, dialect, rows


def analyze_csv_columns(raw: bytes, entities: list) -> list:
    """.csv の列名からエンティティ種別を推定し、確認用の列一覧を返す。"""
    _text, _dialect, rows = _read_csv_rows(raw)
    if not rows:
        return [{"group": None, "columns": []}]

    header = rows[0]
    data_rows = rows[1:21]  # サンプル値の取得は先頭20行程度で十分
    allowed = set(entities or config.ALL_ENTITY_CODES)

    columns = []
    for col_index, col_header in enumerate(header):
        sample = ""
        for row in data_rows:
            if col_index < len(row) and row[col_index]:
                sample = row[col_index]
                break
        columns.append(_column_entry(str(col_index), col_header, allowed, sample))

    return [{"group": None, "columns": columns}]


def mask_csv_file(raw: bytes, entities: list, style: str, column_overrides: dict = None):
    """.csv ファイルの各セルをマスキングする。行・列構造は保持する。

    1行目はヘッダー行とみなし、列名(氏名・住所など)からエンティティ種別を
    推定できた列は、値の内容によらずセル全体をそのエンティティ種別として
    マスキングする(ヘッダー行自体はマスキング対象外)。
    column_overrides を指定した場合、自動判定の代わりにその内容
    (列インデックスの文字列 -> エンティティ種別 または null)を使用する。
    """
    text, dialect, rows = _read_csv_rows(raw)

    all_detections = []
    if not rows:
        return b"", all_detections, text

    header = rows[0]
    allowed = set(entities or config.ALL_ENTITY_CODES)
    if column_overrides is not None:
        forced_types = [column_overrides.get(str(i)) for i in range(len(header))]
    else:
        forced_types = [column_matcher.match_entity_type(h, allowed) for h in header]

    masked_rows = [header]
    for row_index, row in enumerate(rows[1:], start=1):
        masked_row = []
        for col_index, cell in enumerate(row):
            if not cell:
                masked_row.append(cell)
                continue
            forced_type = forced_types[col_index] if col_index < len(forced_types) else None
            masked_cell, detections = _mask_cell(cell, forced_type, entities, style)
            masked_row.append(masked_cell)
            for detection in detections:
                detection["row"] = row_index
                detection["column"] = col_index
                all_detections.append(detection)
        masked_rows.append(masked_row)

    output = io.StringIO()
    writer = csv.writer(output, dialect=dialect)
    writer.writerows(masked_rows)

    return output.getvalue().encode("utf-8-sig"), all_detections, text


def analyze_xlsx_columns(raw: bytes, entities: list) -> list:
    """.xlsx の各シートの列名からエンティティ種別を推定し、確認用の列一覧を返す。"""
    workbook = _open_document(lambda: openpyxl.load_workbook(io.BytesIO(raw)), "Excelファイル")
    allowed = set(entities or config.ALL_ENTITY_CODES)
    groups = []

    for sheet in workbook.worksheets:
        header_row_num = sheet.min_row
        header_cells = list(next(sheet.iter_rows(min_row=header_row_num, max_row=header_row_num), []))

        samples = {}
        for row in sheet.iter_rows(min_row=header_row_num + 1, max_row=header_row_num + 21):
            for cell in row:
                if cell.column not in samples and isinstance(cell.value, str) and cell.value:
                    samples[cell.column] = cell.value

        columns = []
        for cell in header_cells:
            if not isinstance(cell.value, str) or not cell.value:
                continue
            key = f"{sheet.title}:{cell.column}"
            columns.append(_column_entry(key, cell.value, allowed, samples.get(cell.column, "")))

        groups.append({"group": sheet.title, "columns": columns})

    return groups


def mask_xlsx_file(raw: bytes, entities: list, style: str, column_overrides: dict = None):
    """.xlsx ファイルの各セルをマスキングする。シート・書式は保持する。

    各シートの先頭行はヘッダー行とみなし、列名(氏名・住所など)から
    エンティティ種別を推定できた列は、値の内容によらずセル全体をその
    エンティティ種別としてマスキングする(ヘッダー行自体はマスキング対象外)。
    column_overrides を指定した場合、キーは "シート名:列番号" とし、
    自動判定の代わりにその内容を使用する。
    """
    workbook = _open_document(lambda: openpyxl.load_workbook(io.BytesIO(raw)), "Excelファイル")
    all_detections = []
    original_texts = []
    allowed = set(entities or config.ALL_ENTITY_CODES)

    for sheet in workbook.worksheets:
        header_row_num = sheet.min_row
        forced_types = {}
        for cell in next(sheet.iter_rows(min_row=header_row_num, max_row=header_row_num), []):
            if not isinstance(cell.value, str) or not cell.value:
                continue
            if column_overrides is not None:
                forced_types[cell.column] = column_overrides.get(f"{sheet.title}:{cell.column}")
            else:
                forced_types[cell.column] = column_matcher.match_entity_type(cell.value, allowed)

        for row in sheet.iter_rows(min_row=header_row_num + 1):
            for cell in row:
                if not isinstance(cell.value, str) or not cell.value:
                    continue
                original_texts.append(cell.value)
                forced_type = forced_types.get(cell.column)
                masked, detections = _mask_cell(cell.value, forced_type, entities, style)
                if masked != cell.value:
                    cell.value = masked
                for detection in detections:
                    detection["sheet"] = sheet.title
                    detection["row"] = cell.row
                    detection["column"] = cell.column
                all_detections.extend(detections)

    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue(), all_detections, "\n".join(original_texts)


def _mask_paragraph_runs(paragraph, entities: list, style: str, original_texts: list, confirmed: set = None):
    """段落内のテキストをまとめてマスキングし、先頭の run に書き戻す。

    run単位の書式(太字など)は先頭run以外は失われるが、段落の構造・順序は保持する。
    """
    text = paragraph.text
    if not text.strip() or not paragraph.runs:
        return []

    original_texts.append(text)
    results = _analyze_with_label_hints(text, entities)
    masked, detections = mask_resolved(text, results, style, entities, confirmed)
    if masked != text:
        paragraph.runs[0].text = masked
        for run in paragraph.runs[1:]:
            run.text = ""
    return detections


def _iter_docx_texts(document):
    """.docx 内の非空段落テキストを、本文・表・ヘッダー/フッターの順に列挙する(読み取り専用)。"""

    def paragraphs_text(paragraphs):
        for paragraph in paragraphs:
            if paragraph.text.strip():
                yield paragraph.text

    def tables_text(tables):
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    yield from paragraphs_text(cell.paragraphs)
                    yield from tables_text(cell.tables)

    yield from paragraphs_text(document.paragraphs)
    yield from tables_text(document.tables)
    for section in document.sections:
        yield from paragraphs_text(section.header.paragraphs)
        yield from paragraphs_text(section.footer.paragraphs)


def analyze_docx_candidates(raw: bytes, entities: list) -> list:
    document = _open_document(lambda: Document(io.BytesIO(raw)), "Wordファイル")
    return _collect_candidates(_iter_docx_texts(document), entities)


def mask_docx_file(raw: bytes, entities: list, style: str, confirmed: set = None):
    """.docx ファイルの本文・表・ヘッダー/フッターをマスキングする。"""
    document = _open_document(lambda: Document(io.BytesIO(raw)), "Wordファイル")
    all_detections = []
    original_texts = []

    def process_paragraphs(paragraphs):
        for paragraph in paragraphs:
            all_detections.extend(
                _mask_paragraph_runs(paragraph, entities, style, original_texts, confirmed)
            )

    def process_tables(tables):
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    process_paragraphs(cell.paragraphs)
                    process_tables(cell.tables)

    process_paragraphs(document.paragraphs)
    process_tables(document.tables)
    for section in document.sections:
        process_paragraphs(section.header.paragraphs)
        process_paragraphs(section.footer.paragraphs)

    output = io.BytesIO()
    document.save(output)
    return output.getvalue(), all_detections, "\n".join(original_texts)


def _iter_pptx_texts(presentation):
    """.pptx 内の非空段落テキストを、スライド・図形の順に列挙する(読み取り専用)。"""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        yield paragraph.text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            if paragraph.text.strip():
                                yield paragraph.text


def analyze_pptx_candidates(raw: bytes, entities: list) -> list:
    presentation = _open_document(lambda: Presentation(io.BytesIO(raw)), "PowerPointファイル")
    return _collect_candidates(_iter_pptx_texts(presentation), entities)


def mask_pptx_file(raw: bytes, entities: list, style: str, confirmed: set = None):
    """.pptx ファイルのテキストボックス・表をマスキングする。"""
    presentation = _open_document(lambda: Presentation(io.BytesIO(raw)), "PowerPointファイル")
    all_detections = []
    original_texts = []

    def process_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            all_detections.extend(
                _mask_paragraph_runs(paragraph, entities, style, original_texts, confirmed)
            )

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                process_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_text_frame(cell.text_frame)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue(), all_detections, "\n".join(original_texts)


def _open_pdf(raw: bytes):
    document = _open_document(lambda: fitz.open(stream=raw, filetype="pdf"), "PDF")
    if document.needs_pass:
        document.close()
        raise ValueError("パスワードで保護されたPDFは処理できません。")
    return document


# 黒塗り矩形をわずかに広げるマージン(ポイント)。フォントの描画メトリクスの
# 誤差により、検出した矩形ぎりぎりに元の文字の断片が残ることがあるため。
REDACT_PADDING = 1.5


def _pad_rect(rect):
    return fitz.Rect(
        rect.x0 - REDACT_PADDING,
        rect.y0 - REDACT_PADDING,
        rect.x1 + REDACT_PADDING,
        rect.y1 + REDACT_PADDING,
    )


def _redact_pdf_rect(page, rect, style: str, label: str, source_text: str):
    rect = _pad_rect(rect)
    if style == "tag":
        page.add_redact_annot(
            rect, text=f"[{label}]", fontname="japan-s",
            fill=(0, 0, 0), text_color=(1, 1, 1),
        )
    elif style == "mask":
        page.add_redact_annot(
            rect, text="*" * min(len(source_text), 12), fontname="japan-s",
            fill=(0, 0, 0), text_color=(1, 1, 1),
        )
    else:  # redact
        page.add_redact_annot(rect, fill=(0, 0, 0))


def analyze_pdf_candidates(raw: bytes, entities: list) -> list:
    document = _open_pdf(raw)
    try:
        texts = [page.get_text("text") for page in document]
    finally:
        document.close()
    return _collect_candidates(texts, entities)


def mask_pdf_file(raw: bytes, entities: list, style: str, confirmed: set = None):
    """.pdf ファイルをマスキングする。

    PDFはテキストをその場で書き換えると元のレイアウトが崩れるため、検出箇所を
    黒塗り(redaction)することでマスキングする。style に応じて黒塗り部分に
    ラベルやマスク文字を重ねて表示する。confirmed を指定した場合、
    (entity_type, 一致テキスト) が confirmed に含まれる検出のみ黒塗りする。
    """
    document = _open_pdf(raw)
    all_detections = []
    original_texts = []

    for page in document:
        page_text = page.get_text("text")
        if not page_text.strip():
            continue
        original_texts.append(page_text)

        results = _analyze_with_label_hints(page_text, entities)
        seen = set()
        for result in results:
            matched_text = page_text[result.start : result.end].strip()
            if confirmed is not None and (result.entity_type, matched_text) not in confirmed:
                continue
            search_query = re.sub(r"\s+", " ", matched_text).strip()
            if not search_query or search_query in seen:
                continue
            seen.add(search_query)

            label = config.ENTITY_LABELS_JA.get(result.entity_type, result.entity_type)

            rects = page.search_for(search_query)
            if not rects:
                # PDF内部でテキストが改行等により分断され、空白に置き換えた
                # 文字列では一致しない場合、空白・改行区切りの断片ごとに
                # 個別検索してフォールバックする(例: メールアドレスの途中で
                # 改行されているケース)。
                segments = [s for s in re.split(r"\s+", matched_text) if s]
                if len(segments) < 2:
                    continue
                rects = [r for segment in segments for r in page.search_for(segment)]
                if not rects:
                    continue

            # 1つの検出結果が複数の矩形にまたがる場合(改行を挟む場合や、
            # search_for が複数行にわたって一致を返す場合)、タグ/マスク文字の
            # 上書き表示は最初の矩形にのみ行い、残りは黒塗りのみとする
            # (同じラベルが何度も重複表示されるのを避けるため)。
            for idx, rect in enumerate(rects):
                if idx == 0:
                    _redact_pdf_rect(page, rect, style, label, search_query)
                else:
                    page.add_redact_annot(_pad_rect(rect), fill=(0, 0, 0))

            all_detections.append(
                {
                    "entity_type": result.entity_type,
                    "entity_label": label,
                    "score": round(result.score, 3),
                    "text": matched_text,
                    "page": page.number + 1,
                }
            )

        page.apply_redactions()

    output = io.BytesIO()
    document.save(output)
    document.close()
    return output.getvalue(), all_detections, "\n\n".join(original_texts)


def _mask_json_value(value, key, entities, style, allowed, path, all_detections, original_texts, overrides=None):
    if isinstance(value, dict):
        return {
            k: _mask_json_value(
                v, k, entities, style, allowed, f"{path}.{k}", all_detections, original_texts, overrides
            )
            for k, v in value.items()
        }
    if isinstance(value, list):
        return [
            _mask_json_value(
                v, key, entities, style, allowed, f"{path}[{i}]", all_detections, original_texts, overrides
            )
            for i, v in enumerate(value)
        ]
    if isinstance(value, str) and value:
        original_texts.append(value)
        if overrides is not None:
            forced_type = overrides.get(key) if key else None
        else:
            forced_type = column_matcher.match_entity_type(key, allowed) if key else None
        masked, detections = _mask_cell(value, forced_type, entities, style)
        for detection in detections:
            detection["path"] = path
        all_detections.extend(detections)
        return masked
    return value


def _json_records(data):
    return data if isinstance(data, list) else [data]


def analyze_json_columns(raw: bytes, entities: list) -> list:
    """.json のトップレベルのキー名からエンティティ種別を推定し、確認用の一覧を返す。

    (配列内の)レコードごとに異なるキー構成である可能性があるため、
    先頭から最大200レコード分のキーを走査して和集合を取る。
    """
    text, _encoding = decode_bytes(raw)
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"JSONの解析に失敗しました: {exc}")

    allowed = set(entities or config.ALL_ENTITY_CODES)
    samples = {}
    order = []
    for record in _json_records(data)[:200]:
        if not isinstance(record, dict):
            continue
        for k, v in record.items():
            if k not in samples:
                samples[k] = ""
                order.append(k)
            if not samples[k] and isinstance(v, str) and v:
                samples[k] = v

    columns = [_column_entry(key, key, allowed, samples[key]) for key in order]
    return [{"group": None, "columns": columns}]


def mask_json_file(raw: bytes, entities: list, style: str, column_overrides: dict = None):
    """.json ファイルをマスキングする。

    レコードの配列(例: [{"氏名": "...", "住所": "..."}, ...])やネストした
    オブジェクトを再帰的に走査し、文字列の値をマスキングする。オブジェクトの
    キー名(氏名・住所など)からエンティティ種別を推定できる場合は、値の内容に
    よらずその値全体をそのエンティティ種別としてマスキングする。
    column_overrides を指定した場合、キー名(トップレベルのフィールド名)に対する
    自動判定を、その内容(キー名 -> エンティティ種別 または null)で上書きする。
    """
    text, _encoding = decode_bytes(raw)
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"JSONの解析に失敗しました: {exc}")

    allowed = set(entities or config.ALL_ENTITY_CODES)
    all_detections = []
    original_texts = []

    masked_data = _mask_json_value(
        data, None, entities, style, allowed, "$", all_detections, original_texts, column_overrides
    )

    masked_bytes = json.dumps(masked_data, ensure_ascii=False, indent=2).encode("utf-8")
    return masked_bytes, all_detections, "\n".join(original_texts)
